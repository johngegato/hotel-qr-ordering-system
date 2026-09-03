from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Colors
NAVY = RGBColor(2, 6, 23)
SLATE = RGBColor(15, 23, 42)
INDIGO = RGBColor(79, 70, 229)
GOLD = RGBColor(251, 191, 36)
GREEN = RGBColor(52, 211, 153)
WHITE = RGBColor(255, 255, 255)
MUTED = RGBColor(148, 163, 184)
RED = RGBColor(239, 68, 68)


def add_title_slide(title, subtitle=None):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = NAVY

    bar = slide.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(0.22))
    bar.fill.solid()
    bar.fill.fore_color.rgb = INDIGO
    bar.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(11.5), Inches(1.0))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.alignment = PP_ALIGN.LEFT
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE

    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(10.5), Inches(0.8))
        tf2 = sub.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = subtitle
        p2.font.size = Pt(16)
        p2.font.color.rgb = MUTED

    # accent circle
    shape = slide.shapes.add_shape(1, Inches(10.85), Inches(1.15), Inches(1.6), Inches(1.6))
    shape.fill.solid()
    shape.fill.fore_color.rgb = INDIGO
    shape.line.fill.background()
    shape.rotation = 15
    tx = shape.text_frame
    p3 = tx.paragraphs[0]
    p3.text = '🏨'
    p3.alignment = PP_ALIGN.CENTER
    p3.font.size = Pt(26)
    return slide


def add_section_slide(title):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = SLATE

    bar = slide.shapes.add_shape(1, Inches(0.4), Inches(0.5), Inches(12.5), Inches(0.14))
    bar.fill.solid()
    bar.fill.fore_color.rgb = GOLD
    bar.line.fill.background()

    tb = slide.shapes.add_textbox(Inches(0.7), Inches(1.0), Inches(11.5), Inches(1.2))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = WHITE
    return slide


def add_bullets_slide(title, bullets, accent=INDIGO):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = SLATE

    header = slide.shapes.add_shape(1, Inches(0.55), Inches(0.45), Inches(12.1), Inches(0.65))
    header.fill.solid()
    header.fill.fore_color.rgb = accent
    header.line.fill.background()
    tfh = header.text_frame
    p = tfh.paragraphs[0]
    p.text = title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.LEFT

    box = slide.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(11.7), Inches(4.8))
    tf = box.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f'• {b}'
        p.level = 0
        p.space_after = Pt(10)
        p.font.size = Pt(20)
        p.font.color.rgb = WHITE

    return slide


def add_two_column_slide(title, left_title, left_bullets, right_title, right_bullets, accent=INDIGO):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = SLATE

    header = slide.shapes.add_textbox(Inches(0.7), Inches(0.35), Inches(11.7), Inches(0.7))
    tfh = header.text_frame
    p = tfh.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = WHITE

    left = slide.shapes.add_shape(1, Inches(0.8), Inches(1.2), Inches(5.7), Inches(5.2))
    left.fill.solid()
    left.fill.fore_color.rgb = RGBColor(15, 23, 42)
    left.line.color.rgb = accent
    left.line.width = Pt(1.2)
    ltf = left.text_frame
    ltf.word_wrap = True
    p = ltf.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = accent
    for b in left_bullets:
        p = ltf.add_paragraph()
        p.text = f'• {b}'
        p.font.size = Pt(15)
        p.font.color.rgb = WHITE
        p.space_after = Pt(8)

    right = slide.shapes.add_shape(1, Inches(6.9), Inches(1.2), Inches(5.7), Inches(5.2))
    right.fill.solid()
    right.fill.fore_color.rgb = RGBColor(15, 23, 42)
    right.line.color.rgb = GREEN
    right.line.width = Pt(1.2)
    rtf = right.text_frame
    rtf.word_wrap = True
    p = rtf.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = GREEN
    for b in right_bullets:
        p = rtf.add_paragraph()
        p.text = f'• {b}'
        p.font.size = Pt(15)
        p.font.color.rgb = WHITE
        p.space_after = Pt(8)
    return slide


def add_simple_summary(title, metric_lines):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = NAVY

    box = slide.shapes.add_shape(1, Inches(0.7), Inches(0.7), Inches(11.9), Inches(0.8))
    box.fill.solid()
    box.fill.fore_color.rgb = INDIGO
    box.line.fill.background()
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = WHITE

    y = 1.9
    for label, value in metric_lines:
        card = slide.shapes.add_shape(1, Inches(0.9), Inches(y), Inches(3.4), Inches(1.4))
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(30, 41, 59)
        card.line.color.rgb = RGBColor(148, 163, 184)
        tfc = card.text_frame
        p = tfc.paragraphs[0]
        p.text = label
        p.font.size = Pt(12)
        p.font.color.rgb = MUTED
        p2 = tfc.add_paragraph()
        p2.text = value
        p2.font.size = Pt(21)
        p2.font.bold = True
        p2.font.color.rgb = WHITE
        y += 1.8

    return slide


# Slide 1
add_title_slide(
    'Hotel QR Ordering System',
    'Guest Experience + Staff Operations + Admin Control in One Hospitality Platform'
)

# Slide 2
add_bullets_slide(
    'Project Overview',
    [
        'Monorepo solution for a hotel guest experience and staff operations platform.',
        'Built around a QR-based guest flow, real-time staff queues, and centralized admin controls.',
        'Designed for multi-property operation using hotel_id scoping and shared Supabase infrastructure.',
        'Primary goals: speed, transparency, automation, and reduced front desk workload.'
    ],
    accent=INDIGO,
)

# Slide 3
add_two_column_slide(
    'Core Business Value',
    'For Hotel Guests',
    ['Scan room QR code to order food, request services, or book spa treatments.', 'No app download required for the guest experience.', 'Live status tracking reduces uncertainty and repetitive front desk calls.'],
    'For Staff',
    ['Respond to real-time queue items from one tablet screen.', 'View SLA, pending requests, and call handling in one workflow.', 'Reduce manual coordination and improve service response times.'],
    accent=GOLD,
)

# Slide 4
add_bullets_slide(
    'System Architecture',
    [
        'Frontend apps: Next.js guest/admin web portal and Expo React Native staff app.',
        'Shared backend: Supabase PostgreSQL with row-level security, tables, audits, and realtime subscriptions.',
        'Guest QR code flow validates room_id + qr_auth_hash before creating a secure session.',
        'All records are tenant-scoped by hotel_id to support hotel-level separation and data safety.'
    ],
    accent=GREEN,
)

# Slide 5
add_two_column_slide(
    'Guest Experience Flow',
    'Guest Web App',
    ['Room QR opens a guest experience page.', 'Guest can order dining, request housekeeping/tasks, book spa, and call front desk.', 'Request status updates appear in real time through the guest UI.'],
    'Guest Journey',
    ['One-tap actions reduce friction and response time.', 'Session persists through QR validation and guest session tracking.', 'Dynamic branding pulls the hotel name from settings instead of hardcoded values.'],
    accent=INDIGO,
)

# Slide 6
add_bullets_slide(
    'Staff App Workflow',
    [
        'Staff login opens a central operations dashboard with active request queues.',
        'Queues include call requests, spa bookings, food orders, and room tasks.',
        'The app handles status changes, claim actions, SLA checks, and guest phone calls.',
        'It can also send push alerts, reminders, and live call notifications for urgent jobs.'
    ],
    accent=GOLD,
)

# Slide 7
add_section_slide('Admin Web Portal')
# add content in same slide by replacing with actual text boxes? create unique slide below.
admin = prs.slides.add_slide(prs.slide_layouts[5])
admin_bg = admin.background.fill
admin_bg.solid(); admin_bg.fore_color.rgb = SLATE
header = admin.shapes.add_shape(1, Inches(0.5), Inches(0.45), Inches(12.2), Inches(0.7))
header.fill.solid(); header.fill.fore_color.rgb = INDIGO; header.line.fill.background()
htf = header.text_frame; p = htf.paragraphs[0]; p.text = 'Admin Web Portal'; p.font.bold = True; p.font.size = Pt(24); p.font.color.rgb = WHITE
box = admin.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(11.7), Inches(4.8)); tf = box.text_frame; tf.word_wrap = True
for text in [
    '• Hotel settings and branding are centrally managed by property name, phone, theme, and logo.',
    '• Staff account management includes CRUD controls, role assignment, and access status.',
    '• Analytics provide revenue, SLA, and operational performance insights.',
    '• Function room booking module includes rooms, rental equipment, bookings, and scheduling flows.'
]:
    p = tf.paragraphs[0] if not tf.paragraphs[0].text else tf.add_paragraph()
    p.text = text
    p.font.size = Pt(20)
    p.font.color.rgb = WHITE
    p.space_after = Pt(14)

# Slide 8
add_two_column_slide(
    'Key Modules',
    'Guest & Ops',
    ['Food ordering', 'Spa reservations', 'Task requests', 'Guest call handling', 'Live status tracking'],
    'Admin & Backend',
    ['Hotel settings', 'Analytics dashboard', 'Function room booking', 'Audit logs', 'User management'],
    accent=GREEN,
)

# Slide 9
add_bullets_slide(
    'Technology Stack',
    [
        'Next.js for guest-facing and administrative web portals.',
        'Expo / React Native for the staff tablet app on Android.',
        'Supabase for Postgres, auth, row-level security, realtime, and notifications.',
        'Additional integrations: Agora RTC for live voice calls, Expo OTA updates, and push notifications.'
    ],
    accent=INDIGO,
)

# Slide 10
add_bullets_slide(
    'Project Outcome',
    [
        'A single hospitality operations system unifying guest requests, staff queues, and admin controls.',
        'The platform reduces hotel service friction while improving accountability and response time.',
        'It is structured for growth into additional departments, hotel brands, and future automation features.',
        'This repo is suitable for company presentation as a realistic hotel operations platform MVP / production-ready prototype.'
    ],
    accent=GOLD,
)

# Final slide
add_title_slide(
    'Thank You',
    'Questions and discussion: functionality, architecture, deployment, and future expansion options'
)

output_path = r"C:\AMD\New folder\hotel-qr-ordering-system\hotel-qr-ordering-system-overview.pptx"
prs.save(output_path)
print(f"Presentation created: {output_path}")
